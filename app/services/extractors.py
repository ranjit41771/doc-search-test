"""Text extraction pipeline for file documents.

Each extractor returns List[Dict] where each dict is:
    { "page_number": int, "text": str }

All CPU-heavy calls (Tesseract OCR, image rendering) are run in a thread pool
executor via asyncio.get_event_loop().run_in_executor so they do not block the
event loop.
"""

import asyncio
import io
import logging
from functools import partial
from typing import Dict, List

log = logging.getLogger(__name__)

# MIME type → extractor mapping
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    "text/markdown",
    "image/png",
    "image/jpeg",
    "image/tiff",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "text/csv",
}


# ── Sync extraction helpers (run in executor) ──────────────────────────────────

def _ocr_image_bytes(image_bytes: bytes) -> str:
    """Run Tesseract OCR on raw image bytes. Returns extracted text."""
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(image_bytes))
    return pytesseract.image_to_string(img, lang="eng")


def _extract_pdf_sync(file_bytes: bytes) -> List[Dict]:
    """Extract text from PDF per-page.

    Strategy per page:
    1. Use pdfplumber to extract text.
    2. If text is empty/whitespace → render page to image via pdf2image + OCR.
    3. Also OCR any embedded images found on the page.
    """
    import pdfplumber
    import pytesseract
    from pdf2image import convert_from_bytes
    from PIL import Image

    pages = []
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text() or ""

                # If native text is empty/whitespace → scanned page, run OCR
                if not page_text.strip():
                    try:
                        images = convert_from_bytes(
                            file_bytes,
                            first_page=page_num,
                            last_page=page_num,
                            dpi=200,
                        )
                        if images:
                            page_text = pytesseract.image_to_string(images[0], lang="eng")
                    except Exception as e:
                        log.warning("OCR failed for PDF page %d: %s", page_num, e)

                # Extract and OCR embedded images within the page
                try:
                    for img_obj in page.images:
                        raw = img_obj.get("stream")
                        if raw is not None:
                            img_bytes = raw.get_data() if hasattr(raw, "get_data") else raw
                            try:
                                img_text = _ocr_image_bytes(img_bytes)
                                if img_text.strip():
                                    page_text += "\n" + img_text
                            except Exception:
                                pass
                except Exception:
                    pass

                pages.append({"page_number": page_num, "text": page_text.strip()})
    except Exception as e:
        log.error("PDF extraction failed: %s", e)
        raise

    return pages


def _extract_docx_sync(file_bytes: bytes) -> List[Dict]:
    """Extract text from DOCX, including embedded images (OCR'd). Returns single page."""
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(file_bytes))
    text_parts = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    # Embedded images via inline shapes
    for shape in doc.inline_shapes:
        try:
            img_part = shape._inline.graphic.graphicData.pic.blipFill.blip
            rId = img_part.embed
            img_bytes = doc.part.related_parts[rId].blob
            img_text = _ocr_image_bytes(img_bytes)
            if img_text.strip():
                text_parts.append(img_text)
        except Exception:
            pass

    return [{"page_number": 1, "text": "\n".join(text_parts)}]


def _extract_pptx_sync(file_bytes: bytes) -> List[Dict]:
    """Extract text from PPTX per-slide. Each slide = one page."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(io.BytesIO(file_bytes))
    pages = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts = []

        # Text frames
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        text_parts.append(line)

            # Embedded images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    img_text = _ocr_image_bytes(img_bytes)
                    if img_text.strip():
                        text_parts.append(img_text)
                except Exception:
                    pass

        pages.append({"page_number": slide_num, "text": "\n".join(text_parts)})

    return pages


def _extract_txt_sync(file_bytes: bytes) -> List[Dict]:
    """Decode plain text using chardet encoding detection."""
    import chardet

    detected = chardet.detect(file_bytes)
    encoding = detected.get("encoding") or "utf-8"
    try:
        text = file_bytes.decode(encoding)
    except (UnicodeDecodeError, LookupError):
        text = file_bytes.decode("utf-8", errors="replace")
    return [{"page_number": 1, "text": text}]


def _extract_image_sync(file_bytes: bytes) -> List[Dict]:
    """Run Tesseract OCR directly on image bytes. Returns single page."""
    text = _ocr_image_bytes(file_bytes)
    return [{"page_number": 1, "text": text}]


def _extract_xlsx_sync(file_bytes: bytes) -> List[Dict]:
    """Extract cell values from XLSX/CSV. Each sheet = one page."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    pages = []

    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("\t".join(cells))
        pages.append({"page_number": sheet_num, "text": "\n".join(rows)})

    wb.close()
    return pages


def _extract_csv_sync(file_bytes: bytes) -> List[Dict]:
    """Extract CSV as plain text after encoding detection."""
    return _extract_txt_sync(file_bytes)


# ── Async wrappers ─────────────────────────────────────────────────────────────

async def extract_pdf(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_pdf_sync, file_bytes))


async def extract_docx(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_docx_sync, file_bytes))


async def extract_pptx(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_pptx_sync, file_bytes))


async def extract_txt(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_txt_sync, file_bytes))


async def extract_image(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_image_sync, file_bytes))


async def extract_xlsx(file_bytes: bytes) -> List[Dict]:
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, partial(_extract_xlsx_sync, file_bytes))


# ── Router ─────────────────────────────────────────────────────────────────────

async def route_extraction(file_bytes: bytes, mime_type: str) -> List[Dict]:
    """Dispatch file bytes to the correct extractor based on MIME type."""
    mime_type = mime_type.lower().split(";")[0].strip()

    if mime_type == "application/pdf":
        return await extract_pdf(file_bytes)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return await extract_docx(file_bytes)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return await extract_pptx(file_bytes)
    elif mime_type in ("text/plain", "text/markdown"):
        return await extract_txt(file_bytes)
    elif mime_type in ("image/png", "image/jpeg", "image/tiff"):
        return await extract_image(file_bytes)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return await extract_xlsx(file_bytes)
    elif mime_type == "text/csv":
        return await extract_txt(file_bytes)
    else:
        raise ValueError(f"Unsupported MIME type: {mime_type}")


# ── Chunking ───────────────────────────────────────────────────────────────────

def chunk_pages(
    pages: List[Dict],
    chunk_size: int = 500,
    overlap: int = 100,
) -> List[Dict]:
    """Split page texts into overlapping chunks for ES indexing.

    Returns list of:
        { "chunk_index": int, "page_number": int, "text": str }

    Overlap carries the last `overlap` characters of each chunk into the
    beginning of the next chunk so that term matches near chunk boundaries
    are not missed.
    """
    chunks = []
    chunk_index = 0

    for page in pages:
        text = page["text"].strip()
        page_number = page["page_number"]
        if not text:
            continue

        start = 0
        while start < len(text):
            end = start + chunk_size
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "chunk_index": chunk_index,
                    "page_number": page_number,
                    "text": chunk_text,
                })
                chunk_index += 1
            if end >= len(text):
                break
            # Move forward by (chunk_size - overlap) to create overlap window
            start = end - overlap

    return chunks
